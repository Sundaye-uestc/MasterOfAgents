#!/usr/bin/env python3
"""
PPTX to HTML Preview Converter

Converts a .pptx file into a standalone HTML viewer with slide navigation.
Extracts text, images, and basic layout from the PPTX and renders slides as
HTML/CSS pages. The output can be served as a webpage artifact in AgentHub.

Usage:
  python pptx_to_preview.py --input presentation.pptx --output preview_dir
"""

import argparse
import json
import os
import re
import shutil
import sys
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from xml.etree import ElementTree as ET


# =============================================================================
# Constants
# =============================================================================

PPT_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAWING_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NAMESPACE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

EMU_PER_INCH = 914400
DEFAULT_SLIDE_W = 12192000  # 10 inches in EMU
DEFAULT_SLIDE_H = 6858000   # 7.5 inches in EMU


# =============================================================================
# PPTX Parsing
# =============================================================================

def _qname(tag: str, ns: str = PPT_NAMESPACE) -> str:
    """Build a qualified XML tag name."""
    return f"{{{ns}}}{tag}"


def _d(tag: str) -> str:
    """DrawingML qualified name."""
    return _qname(tag, DRAWING_NAMESPACE)


def _clean_text(text: str) -> str:
    """Clean text for HTML output."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _emu_to_px(emu: int, dpi: int = 96) -> float:
    """Convert EMU (English Metric Units) to pixels."""
    return emu / EMU_PER_INCH * dpi


def _parse_color(element: ET.Element, attr: str) -> Optional[str]:
    """Extract color from an XML element's sRGB child."""
    srgb = element.find(_d("srgbClr"))
    if srgb is not None:
        return srgb.get("val", "")
    return None


def extract_slide_texts(pptx_path: Path) -> List[List[Dict[str, Any]]]:
    """
    Extract text content and position from each slide.
    Returns a list of slides, each a list of text elements with position info.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx not installed", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width or DEFAULT_SLIDE_W
    slide_height = prs.slide_height or DEFAULT_SLIDE_H

    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            shape_data = {
                "left": _emu_to_px(shape.left or 0),
                "top": _emu_to_px(shape.top or 0),
                "width": _emu_to_px(shape.width or 0),
                "height": _emu_to_px(shape.height or 0),
                "shape_type": str(shape.shape_type),
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    para_text = ""
                    runs_data = []
                    for run in para.runs:
                        run_info = {
                            "text": run.text,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "size": run.font.size,
                            "color": str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None,
                        }
                        runs_data.append(run_info)
                        para_text += run.text
                    if para_text.strip():
                        paragraphs.append({
                            "text": para_text.strip(),
                            "runs": runs_data,
                            "alignment": str(para.alignment) if para.alignment else None,
                        })
                if paragraphs:
                    shape_data["paragraphs"] = paragraphs

            if shape.has_table:
                table_data = {"rows": []}
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip() if cell.text else "")
                    table_data["rows"].append(row_data)
                shape_data["table"] = table_data

            # Check for image
            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                shape_data["is_image"] = True
                # Try to get image info
                try:
                    shape_data["image_content_type"] = shape.image.content_type if hasattr(shape, 'image') else None
                except Exception:
                    pass

            texts.append(shape_data)

        slides_data.append({
            "index": slide_idx,
            "shapes": texts,
            "width_px": _emu_to_px(slide_width),
            "height_px": _emu_to_px(slide_height),
            "layout": slide.slide_layout.name if slide.slide_layout else "Custom",
        })

    return slides_data


def extract_images_from_pptx(pptx_path: Path, output_dir: str) -> List[str]:
    """Extract embedded images from PPTX ZIP and save to output/images/."""
    image_dir = os.path.join(output_dir, "images")
    os.makedirs(image_dir, exist_ok=True)

    extracted = []
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            for name in zf.namelist():
                if name.startswith("ppt/media/") and not name.endswith("/"):
                    # Extract and copy to images/
                    basename = os.path.basename(name)
                    dest = os.path.join(image_dir, basename)
                    with zf.open(name) as src:
                        with open(dest, "wb") as dst:
                            dst.write(src.read())
                    extracted.append(basename)
    except Exception as e:
        print(f"Warning: Could not extract images: {e}")

    return extracted


# =============================================================================
# HTML Generation
# =============================================================================

def _rgb_hex_to_css(rgb_str: Optional[str]) -> Optional[str]:
    """Convert python-pptx RGB string to CSS hex color."""
    if not rgb_str:
        return None
    if rgb_str.startswith("#"):
        return rgb_str
    return f"#{rgb_str}"


def _size_to_css(pt_size) -> Optional[str]:
    """Convert font point size to CSS."""
    if pt_size is None:
        return None
    # pt_size from python-pptx is in EMU (1/12700 of a point)
    # Actually it's stored as a float representing points
    try:
        pts = float(pt_size) / 12700  # EMU to points
    except (TypeError, ValueError):
        try:
            pts = float(pt_size)
        except (TypeError, ValueError):
            return None
    return f"{pts:.1f}pt"


def generate_slide_html(slide_data: Dict[str, Any], slide_num: int, total: int) -> str:
    """Generate HTML for a single slide as a positioned div."""
    shapes_html = []
    slide_w = slide_data["width_px"]
    slide_h = slide_data["height_px"]

    for shape in slide_data.get("shapes", []):
        style_parts = [
            f"position:absolute",
            f"left:{shape['left']}px",
            f"top:{shape['top']}px",
            f"width:{shape['width']}px",
            f"height:{shape['height']}px",
            "overflow:hidden",
            "box-sizing:border-box",
            "padding:4px",
        ]

        inner_html = ""

        # Table rendering
        if "table" in shape:
            rows_html = []
            for row in shape["table"]["rows"]:
                cells = "".join(
                    f'<td style="border:1px solid #ddd;padding:4px 8px;">{_clean_text(c)}</td>'
                    for c in row
                )
                rows_html.append(f"<tr>{cells}</tr>")
            inner_html = (
                f'<table style="width:100%;border-collapse:collapse;font-size:11px;">'
                f"{''.join(rows_html)}</table>"
            )

        # Image
        elif shape.get("is_image"):
            style_parts.append("display:flex;align-items:center;justify-content:center;")
            inner_html = '<span style="color:#999;font-size:10px;">[Image]</span>'

        # Text content
        elif "paragraphs" in shape:
            for para in shape["paragraphs"]:
                if not para.get("runs"):
                    para_html = _clean_text(para["text"])
                    inner_html += f'<p style="margin:0 0 4px 0;">{para_html}</p>'
                else:
                    runs_html = []
                    for run in para["runs"]:
                        rtext = _clean_text(run["text"])
                        rstyle = []
                        if run.get("bold"):
                            rstyle.append("font-weight:bold")
                        if run.get("italic"):
                            rstyle.append("font-style:italic")
                        color = _rgb_hex_to_css(run.get("color"))
                        if color:
                            rstyle.append(f"color:{color}")
                        size = _size_to_css(run.get("size"))
                        if size:
                            rstyle.append(f"font-size:{size}")
                        style_attr = ";".join(rstyle)
                        if style_attr:
                            runs_html.append(f'<span style="{style_attr}">{rtext}</span>')
                        else:
                            runs_html.append(rtext)
                    inner_html += f'<p style="margin:0 0 4px 0;">{"".join(runs_html)}</p>'

        style = ";".join(style_parts)
        shapes_html.append(f'<div style="{style}">{inner_html}</div>')

    # Scale slide to fit viewport while maintaining aspect ratio
    aspect = slide_w / max(slide_h, 1)
    target_w = 960
    target_h = int(target_w / aspect)

    return f"""
    <div class="slide" data-slide="{slide_num}">
        <div style="
            position:relative;
            width:{target_w}px;
            height:{target_h}px;
            background:#fff;
            box-shadow:0 2px 20px rgba(0,0,0,0.3);
            margin:0 auto;
            font-family:'Segoe UI',Arial,sans-serif;
            color:#333;
        ">
            {''.join(shapes_html)}
        </div>
        <div class="slide-num">{slide_num} / {total}</div>
    </div>"""


def generate_viewer_html(slides_data: List[Dict[str, Any]], output_dir: str) -> str:
    """Generate the full HTML viewer with slide navigation."""
    total = len(slides_data)
    slide_pages = "\n".join(
        generate_slide_html(s, i + 1, total) for i, s in enumerate(slides_data)
    )

    # Thumbnail navigation dots
    dots = "\n".join(
        f'<span class="dot" data-dot="{i+1}" onclick="goToSlide({i+1})"></span>'
        for i in range(total)
    )

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PPT 预览</title>
    <style>
        * {{ margin:0; padding:0; box-sizing:border-box; }}
        body {{
            background:#1a1a2e;
            font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,sans-serif;
            display:flex;
            flex-direction:column;
            height:100vh;
            overflow:hidden;
        }}
        .viewport {{
            flex:1;
            display:flex;
            align-items:center;
            justify-content:center;
            padding:20px;
            overflow:hidden;
        }}
        .slide {{
            display:none;
            text-align:center;
            animation:fadeIn 0.3s ease;
        }}
        .slide.active {{ display:block; }}
        @keyframes fadeIn {{
            from {{ opacity:0; transform:translateY(8px); }}
            to {{ opacity:1; transform:translateY(0); }}
        }}
        .slide-num {{
            color:#888;
            font-size:11px;
            margin-top:6px;
        }}
        .controls {{
            display:flex;
            align-items:center;
            justify-content:center;
            gap:12px;
            padding:10px 20px;
            background:rgba(255,255,255,0.04);
            border-top:1px solid rgba(255,255,255,0.08);
        }}
        .controls button {{
            background:rgba(255,255,255,0.1);
            border:1px solid rgba(255,255,255,0.15);
            color:#ccc;
            padding:6px 14px;
            border-radius:6px;
            cursor:pointer;
            font-size:13px;
            transition:all 0.15s;
        }}
        .controls button:hover {{
            background:rgba(255,255,255,0.2);
            color:#fff;
        }}
        .controls button:disabled {{
            opacity:0.3;
            cursor:default;
        }}
        .dots {{
            display:flex;
            gap:6px;
        }}
        .dot {{
            width:8px; height:8px;
            border-radius:50%;
            background:rgba(255,255,255,0.2);
            cursor:pointer;
            transition:all 0.15s;
        }}
        .dot.active, .dot:hover {{
            background:rgba(255,255,255,0.7);
        }}
        .page-indicator {{
            color:#aaa;
            font-size:13px;
            min-width:60px;
            text-align:center;
        }}
    </style>
</head>
<body>
    <div class="viewport" id="viewport">
        {slide_pages}
    </div>
    <div class="controls">
        <button onclick="goToSlide(1)" title="首页">⏮</button>
        <button onclick="prevSlide()" title="上一页">◀</button>
        <span class="page-indicator" id="indicator">1 / {total}</span>
        <button onclick="nextSlide()" title="下一页">▶</button>
        <button onclick="goToSlide({total})" title="末页">⏭</button>
        <span style="color:#555;margin:0 8px;">|</span>
        <div class="dots" id="dots">{dots}</div>
    </div>

    <script>
        let current = 1;
        const total = {total};

        function showSlide(n) {{
            document.querySelectorAll('.slide').forEach(s => s.classList.remove('active'));
            const target = document.querySelector(`.slide[data-slide="${{n}}"]`);
            if (target) target.classList.add('active');
            document.querySelectorAll('.dot').forEach(d => d.classList.remove('active'));
            const dot = document.querySelector(`.dot[data-dot="${{n}}"]`);
            if (dot) dot.classList.add('active');
            document.getElementById('indicator').textContent = `${{n}} / ${{total}}`;
        }}

        function nextSlide() {{ if (current < total) {{ current++; showSlide(current); }} }}
        function prevSlide() {{ if (current > 1) {{ current--; showSlide(current); }} }}
        function goToSlide(n) {{ current = Math.max(1, Math.min(total, n)); showSlide(current); }}

        // Keyboard navigation
        document.addEventListener('keydown', e => {{
            if (e.key === 'ArrowRight' || e.key === 'ArrowDown') {{ e.preventDefault(); nextSlide(); }}
            if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') {{ e.preventDefault(); prevSlide(); }}
            if (e.key === 'Home') {{ e.preventDefault(); goToSlide(1); }}
            if (e.key === 'End') {{ e.preventDefault(); goToSlide(total); }}
        }});

        // Touch swipe
        let touchStartX = 0;
        document.addEventListener('touchstart', e => {{ touchStartX = e.touches[0].clientX; }});
        document.addEventListener('touchend', e => {{
            const dx = e.changedTouches[0].clientX - touchStartX;
            if (Math.abs(dx) > 50) {{ dx < 0 ? nextSlide() : prevSlide(); }}
        }});

        // Click on viewport edges
        document.getElementById('viewport').addEventListener('click', e => {{
            const rect = e.currentTarget.getBoundingClientRect();
            if (e.clientX - rect.left < rect.width * 0.2) prevSlide();
            if (e.clientX - rect.left > rect.width * 0.8) nextSlide();
        }});

        showSlide(1);
    </script>
</body>
</html>"""

    html_path = os.path.join(output_dir, "index.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html)
    return html_path


# =============================================================================
# Main Entry Point
# =============================================================================

def create_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Convert PPTX file to HTML preview viewer",
    )
    parser.add_argument(
        "--input", required=True,
        help="Path to the .pptx file",
    )
    parser.add_argument(
        "--output", required=True,
        help="Output directory for the preview viewer",
    )
    return parser


def main() -> None:
    args = create_argument_parser().parse_args()

    pptx_path = Path(args.input)
    if not pptx_path.exists():
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    output_dir = args.output
    os.makedirs(output_dir, exist_ok=True)

    print(f"Converting: {pptx_path.name}")
    print(f"Output: {output_dir}")

    # Extract slide content
    slides_data = extract_slide_texts(pptx_path)

    # Extract embedded images
    images = extract_images_from_pptx(pptx_path, output_dir)

    # Generate viewer
    viewer_path = generate_viewer_html(slides_data, output_dir)

    print(f"Slides: {len(slides_data)}")
    print(f"Images extracted: {len(images)}")
    print(f"Viewer: {viewer_path}")

    # Output structured result
    result = {
        "status": "success",
        "slides": len(slides_data),
        "images_extracted": len(images),
        "viewer_html": viewer_path,
        "output_dir": output_dir,
    }
    print("---RESULT---")
    print(json.dumps(result, ensure_ascii=False, indent=2))
    print("---END RESULT---")


if __name__ == "__main__":
    main()
